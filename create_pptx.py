from pptx import Presentation

# Crear una nueva presentación
prs = Presentation()

# Agregar una diapositiva de título
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]

title.text = "Presentación de Ejemplo"
subtitle.text = "Creada con python-pptx"

# Agregar una diapositiva de contenido con viñetas
slide_content = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_content.shapes.title
content = slide_content.placeholders[1]

title.text = "Contenido de Ejemplo"
content.text = "Aquí están los puntos principales:"
for bullet in ["Punto 1", "Punto 2", "Punto 3"]:
    p = content.text_frame.add_paragraph()
    p.text = bullet

# Agregar una diapositiva con una imagen
slide_image = prs.slides.add_slide(prs.slide_layouts[5])
title = slide_image.shapes.title
title.text = "Diapositiva con Imagen"

# Ruta de la imagen
image_path = "/home/jona/PycharmProjects/Word_pandas/imagen.jpg"
# Ajustar coordenadas y tamaño
left, top, width, height = 1, 1, 4, 4  # en pulgadas
try:
    slide_image.shapes.add_picture(image_path, left * 914400, top * 914400, width * 914400, height * 914400)
except FileNotFoundError:
    print(f"No se encontró la imagen en la ruta: {image_path}")

# Guardar la presentación
output_path = "presentacion_ejemplo.pptx"
prs.save(output_path)
print(f"Presentación guardada como {output_path}")
